import sys
from pptx import Presentation
from pathlib import Path

if len(sys.argv) < 2:
    print('Usage: python analyze_srs.py <path_to_pptx>')
    sys.exit(1)

ppt_path = Path(sys.argv[1])
if not ppt_path.exists():
    print(f'File not found: {ppt_path}')
    sys.exit(1)

prs = Presentation(str(ppt_path))
print(f'SRS Slides: {len(prs.slides)}')
print('--- BEGIN SLIDE TEXT EXTRACT (writing utf-8 to srs_extract.txt) ---')
out_path = Path('srs_extract.txt')
out = out_path.open('w', encoding='utf-8')
for i, slide in enumerate(prs.slides, start=1):
    header = f"\n=== Slide {i} ===\n"
    out.write(header)
    print(f"Slide {i} ...")
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            texts.append(shape.text.strip())
        elif getattr(shape, 'has_text_frame', False):
            buf = []
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    buf.append(p.text.strip())
            if buf:
                texts.append(' '.join(buf))
    if not texts:
        out.write('[No textual content]\n')
    else:
        for t in texts:
            # Collapse excessive whitespace
            out.write('- ' + ' '.join(t.split()) + "\n")
out.write('\n--- END SLIDE TEXT EXTRACT ---\n')
out.close()
print(f'Wrote: {out_path.resolve()}')
